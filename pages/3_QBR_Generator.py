"""
QBR Generator
-------------
Upload a FRDM supplier risk report CSV to generate a PowerPoint
Quarterly Business Review deck for a customer.
"""

import io
from collections import Counter
from datetime import date

import pandas as pd
import streamlit as st
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Page config ───────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="QBR Generator · FRDM Tools",
    page_icon="📋",
    layout="wide",
)

# ── Brand colors ──────────────────────────────────────────────────────────────
FRDM_DARK   = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy
FRDM_BLUE   = RGBColor(0x16, 0x21, 0x3E)   # mid navy
FRDM_ACCENT = RGBColor(0x0F, 0x3C, 0x96)   # brand blue
FRDM_GREEN  = RGBColor(0x28, 0xA7, 0x45)
FRDM_AMBER  = RGBColor(0xFF, 0xA5, 0x00)
FRDM_RED    = RGBColor(0xDC, 0x35, 0x45)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xF4, 0xF6, 0xFB)
MID_GREY    = RGBColor(0x88, 0x8E, 0xA4)

RISK_COLOR_MAP = {
    "Low":          RGBColor(0x28, 0xA7, 0x45),
    "Low Moderate": RGBColor(0x8B, 0xC3, 0x4A),
    "Moderate":     RGBColor(0xFF, 0xA5, 0x00),
    "High":         RGBColor(0xFF, 0x57, 0x22),
    "Elevated":     RGBColor(0xDC, 0x35, 0x45),
    "N/A":          MID_GREY,
}

TRANS_LABELS = ["Very Low", "Low", "Moderate", "High", "Very High"]

# ── Score → label helpers ─────────────────────────────────────────────────────

def risk_label(score):
    if score is None:
        return "N/A"
    if score < 0.2:
        return "Low"
    if score < 0.4:
        return "Low Moderate"
    if score < 0.6:
        return "Moderate"
    if score < 0.8:
        return "High"
    return "Elevated"


def transparency_label(score):
    if score is None:
        return "N/A"
    if score < 0.2:
        return "Very Low"
    if score < 0.4:
        return "Low"
    if score < 0.6:
        return "Moderate"
    if score < 0.8:
        return "High"
    return "Very High"


def parse_pipe(val):
    if not val or (isinstance(val, float)):
        return []
    return [v.strip() for v in str(val).split("|") if v.strip()]


def safe_float(val):
    try:
        return float(val)
    except (TypeError, ValueError):
        return None


# ── PPTX helpers ──────────────────────────────────────────────────────────────

def set_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=12, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def section_header(slide, title: str):
    """Dark accent bar across top with section title."""
    add_rect(slide, Inches(0), Inches(0), Inches(10), Inches(0.7), FRDM_ACCENT)
    add_textbox(slide, title,
                Inches(0.3), Inches(0.08), Inches(9.4), Inches(0.55),
                font_size=18, bold=True, color=WHITE)


def stat_box(slide, label, value, left, top, width=Inches(2.1), height=Inches(1.0),
             value_color=WHITE, bg=FRDM_BLUE):
    add_rect(slide, left, top, width, height, bg)
    add_textbox(slide, str(value),
                left + Inches(0.1), top + Inches(0.05), width - Inches(0.2), Inches(0.55),
                font_size=22, bold=True, color=value_color, align=PP_ALIGN.CENTER)
    add_textbox(slide, label,
                left + Inches(0.05), top + Inches(0.58), width - Inches(0.1), Inches(0.38),
                font_size=9, bold=False, color=MID_GREY, align=PP_ALIGN.CENTER)


def bullet_list(slide, items, left, top, width, height, font_size=10, color=FRDM_DARK, max_items=10):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items[:max_items]):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(3)


def mini_bar(slide, label, count, total, left, top, bar_width=Inches(3.5), color=FRDM_ACCENT):
    """Single horizontal bar with label and count."""
    row_h = Inches(0.28)
    pct = count / total if total else 0
    filled = bar_width * pct

    add_textbox(slide, label,
                left, top, Inches(2.3), row_h,
                font_size=9, color=FRDM_DARK)
    # background track
    add_rect(slide, left + Inches(2.35), top + Inches(0.06),
             bar_width, Inches(0.16), LIGHT_GREY)
    # filled portion
    if filled > Inches(0.05):
        add_rect(slide, left + Inches(2.35), top + Inches(0.06),
                 filled, Inches(0.16), color)
    add_textbox(slide, str(count),
                left + Inches(2.35) + bar_width + Inches(0.05), top,
                Inches(0.4), row_h, font_size=9, color=MID_GREY)
    return top + row_h + Inches(0.04)


# ── Slide builders ────────────────────────────────────────────────────────────

def slide_title(prs, company_name: str, quarter: str):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    set_bg(slide, FRDM_DARK)

    # accent stripe
    add_rect(slide, Inches(0), Inches(2.8), Inches(10), Inches(0.06), FRDM_ACCENT)

    add_textbox(slide, "FRDM",
                Inches(0.6), Inches(0.5), Inches(4), Inches(0.7),
                font_size=14, bold=True, color=FRDM_ACCENT)
    add_textbox(slide, "Quarterly Business Review",
                Inches(0.6), Inches(1.15), Inches(8.8), Inches(0.9),
                font_size=32, bold=True, color=WHITE)
    add_textbox(slide, company_name,
                Inches(0.6), Inches(2.0), Inches(8.8), Inches(0.7),
                font_size=22, bold=False, color=RGBColor(0xCC, 0xD6, 0xFF))
    add_textbox(slide, quarter,
                Inches(0.6), Inches(3.0), Inches(4), Inches(0.4),
                font_size=12, color=MID_GREY)
    add_textbox(slide, "Prepared by FRDM  ·  Confidential",
                Inches(0.6), Inches(6.8), Inches(8.8), Inches(0.4),
                font_size=9, color=MID_GREY)


def slide_executive_summary(prs, stats: dict):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_bg(slide, WHITE)
    section_header(slide, "Executive Summary")

    total = stats["total"]
    avg_risk = stats["avg_risk"]
    avg_trans = stats["avg_trans"]
    risk_dist = stats["risk_dist"]

    # Top stat boxes
    stat_box(slide, "Total Suppliers Monitored", f"{total:,}",
             Inches(0.3), Inches(0.85), bg=FRDM_ACCENT)
    stat_box(slide, "Avg Risk Score", f"{avg_risk:.2f}",
             Inches(2.55), Inches(0.85),
             value_color=RISK_COLOR_MAP.get(risk_label(avg_risk), WHITE), bg=FRDM_BLUE)
    stat_box(slide, "Avg Transparency", f"{avg_trans:.2f}",
             Inches(4.8), Inches(0.85), bg=FRDM_BLUE)
    stat_box(slide, "Suppliers w/ Alerts", str(stats["alert_count"]),
             Inches(7.05), Inches(0.85), bg=FRDM_BLUE)

    # Risk distribution bars
    add_textbox(slide, "Risk Distribution",
                Inches(0.3), Inches(2.1), Inches(4.5), Inches(0.35),
                font_size=12, bold=True, color=FRDM_DARK)

    y = Inches(2.55)
    for label in ["Low", "Low Moderate", "Moderate", "High", "Elevated"]:
        cnt = risk_dist.get(label, 0)
        y = mini_bar(slide, label, cnt, total, Inches(0.3), y,
                     color=RISK_COLOR_MAP[label])

    # Transparency distribution bars
    add_textbox(slide, "Transparency Distribution",
                Inches(5.4), Inches(2.1), Inches(4.5), Inches(0.35),
                font_size=12, bold=True, color=FRDM_DARK)

    y2 = Inches(2.55)
    trans_dist = stats["trans_dist"]
    for label in TRANS_LABELS:
        cnt = trans_dist.get(label, 0)
        y2 = mini_bar(slide, label, cnt, total, Inches(5.4), y2,
                      color=FRDM_ACCENT)

    # Key metrics row at bottom
    add_textbox(slide,
                f"Suppliers flagged with warnings: {stats['warning_count']}   ·   "
                f"Suppliers with documents on file: {stats['doc_count']}   ·   "
                f"Unique regulations tracked: {stats['reg_count']}",
                Inches(0.3), Inches(6.7), Inches(9.4), Inches(0.4),
                font_size=9, color=MID_GREY, align=PP_ALIGN.CENTER)


def slide_risk_landscape(prs, stats: dict, top_risky: list):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_bg(slide, WHITE)
    section_header(slide, "Supply Chain Risk Landscape")

    # Country exposure
    add_textbox(slide, "Top Countries by Supplier Exposure",
                Inches(0.3), Inches(0.85), Inches(4.5), Inches(0.35),
                font_size=12, bold=True, color=FRDM_DARK)
    y = Inches(1.25)
    total = stats["total"]
    for country, cnt in stats["top_countries"][:8]:
        y = mini_bar(slide, country, cnt, total, Inches(0.3), y)

    # Top risky suppliers
    add_textbox(slide, "Highest Risk Suppliers",
                Inches(5.3), Inches(0.85), Inches(4.5), Inches(0.35),
                font_size=12, bold=True, color=FRDM_DARK)

    y2 = Inches(1.25)
    for row in top_risky[:7]:
        score = safe_float(row.get("topline_risk_score"))
        label = risk_label(score)
        color = RISK_COLOR_MAP.get(label, MID_GREY)
        name = str(row.get("Supplier", ""))[:35]
        score_str = f"{score:.2f}" if score is not None else "N/A"
        add_rect(slide, Inches(5.3), y, Inches(4.5), Inches(0.32), LIGHT_GREY)
        add_textbox(slide, name,
                    Inches(5.4), y + Inches(0.04), Inches(3.1), Inches(0.28),
                    font_size=9, color=FRDM_DARK)
        add_rect(slide, Inches(8.9), y, Inches(0.85), Inches(0.32), color)
        add_textbox(slide, f"{score_str} · {label}",
                    Inches(8.92), y + Inches(0.05), Inches(1.2), Inches(0.25),
                    font_size=7, bold=True, color=WHITE)
        y2 = y + Inches(0.36)
        y = y2


def slide_warnings(prs, stats: dict, warning_suppliers: list):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_bg(slide, WHITE)
    section_header(slide, "Warnings & Trade Risk Signals")

    total_warn = stats["warning_count"]
    total = stats["total"]
    top_warnings = stats["top_warnings"]

    stat_box(slide, "Suppliers with Warnings", str(total_warn),
             Inches(0.3), Inches(0.85),
             value_color=FRDM_RED if total_warn > 0 else FRDM_GREEN, bg=FRDM_BLUE)
    pct = round(total_warn / total * 100) if total else 0
    stat_box(slide, "% of Portfolio", f"{pct}%",
             Inches(2.55), Inches(0.85), bg=FRDM_BLUE)

    add_textbox(slide, "Most Common Warning Types",
                Inches(0.3), Inches(2.1), Inches(4.5), Inches(0.35),
                font_size=12, bold=True, color=FRDM_DARK)
    y = Inches(2.5)
    for warn, cnt in top_warnings[:8]:
        y = mini_bar(slide, warn[:40], cnt, total_warn, Inches(0.3), y,
                     color=FRDM_RED)

    # Spotlight suppliers
    add_textbox(slide, "Spotlight: Suppliers with High-Risk Warnings",
                Inches(5.3), Inches(2.1), Inches(4.5), Inches(0.35),
                font_size=12, bold=True, color=FRDM_DARK)
    y2 = Inches(2.5)
    for row in warning_suppliers[:6]:
        name = str(row.get("Supplier", ""))[:32]
        warnings = parse_pipe(row.get("Warnings", ""))
        w_short = warnings[0][:35] if warnings else ""
        score = safe_float(row.get("topline_risk_score"))
        add_rect(slide, Inches(5.3), y2, Inches(4.5), Inches(0.44), LIGHT_GREY)
        add_textbox(slide, name,
                    Inches(5.4), y2 + Inches(0.02), Inches(4.2), Inches(0.22),
                    font_size=9, bold=True, color=FRDM_DARK)
        add_textbox(slide, w_short,
                    Inches(5.4), y2 + Inches(0.22), Inches(4.2), Inches(0.2),
                    font_size=8, color=MID_GREY)
        y2 += Inches(0.48)


def slide_regulations(prs, stats: dict):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_bg(slide, WHITE)
    section_header(slide, "Regulatory Exposure")

    total = stats["total"]

    add_textbox(slide, "Industry Regulations",
                Inches(0.3), Inches(0.85), Inches(4.5), Inches(0.35),
                font_size=12, bold=True, color=FRDM_DARK)
    y = Inches(1.25)
    for reg, cnt in stats["top_ind_regs"][:9]:
        y = mini_bar(slide, reg[:42], cnt, total, Inches(0.3), y,
                     color=FRDM_ACCENT)

    add_textbox(slide, "Product Regulations",
                Inches(5.3), Inches(0.85), Inches(4.5), Inches(0.35),
                font_size=12, bold=True, color=FRDM_DARK)
    y2 = Inches(1.25)
    for reg, cnt in stats["top_prod_regs"][:9]:
        y2 = mini_bar(slide, reg[:42], cnt, total, Inches(5.3), y2,
                      color=FRDM_BLUE)

    add_textbox(slide,
                "Regulations are tracked across your supply chain by FRDM and updated as legislation changes.",
                Inches(0.3), Inches(6.65), Inches(9.4), Inches(0.45),
                font_size=9, color=MID_GREY, align=PP_ALIGN.CENTER)


def slide_documents(prs, stats: dict):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_bg(slide, WHITE)
    section_header(slide, "Supplier Documents & Transparency")

    total = stats["total"]
    doc_count = stats["doc_count"]
    avg_trans = stats["avg_trans"]

    stat_box(slide, "Suppliers with Documents", str(doc_count),
             Inches(0.3), Inches(0.85), bg=FRDM_ACCENT)
    stat_box(slide, "Without Documents", str(total - doc_count),
             Inches(2.55), Inches(0.85), bg=FRDM_BLUE)
    stat_box(slide, "Avg Transparency Score", f"{avg_trans:.2f}",
             Inches(4.8), Inches(0.85), bg=FRDM_BLUE)
    stat_box(slide, "Transparency Level", transparency_label(avg_trans),
             Inches(7.05), Inches(0.85), bg=FRDM_BLUE)

    add_textbox(slide, "Documents on File",
                Inches(0.3), Inches(2.1), Inches(4.5), Inches(0.35),
                font_size=12, bold=True, color=FRDM_DARK)
    y = Inches(2.5)
    for doc, cnt in stats["top_docs"][:8]:
        y = mini_bar(slide, doc[:42], cnt, doc_count or 1, Inches(0.3), y,
                     color=FRDM_GREEN)

    add_textbox(slide,
                "Transparency Insight",
                Inches(5.3), Inches(2.1), Inches(4.5), Inches(0.35),
                font_size=12, bold=True, color=FRDM_DARK)

    trans_dist = stats["trans_dist"]
    y2 = Inches(2.5)
    for label in TRANS_LABELS:
        cnt = trans_dist.get(label, 0)
        y2 = mini_bar(slide, label, cnt, total, Inches(5.3), y2,
                      color=FRDM_GREEN)

    add_textbox(slide,
                "Transparency scores reflect how much verifiable information exists about a supplier's "
                "human rights and labor practices. Higher scores indicate greater disclosure.",
                Inches(0.3), Inches(6.45), Inches(9.4), Inches(0.6),
                font_size=9, color=MID_GREY, align=PP_ALIGN.CENTER)


def slide_frdm_value(prs, stats: dict, company_name: str):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_bg(slide, FRDM_DARK)

    add_rect(slide, Inches(0), Inches(0), Inches(10), Inches(0.06), FRDM_ACCENT)

    add_textbox(slide, "What FRDM Is Doing For You",
                Inches(0.5), Inches(0.25), Inches(9), Inches(0.65),
                font_size=20, bold=True, color=WHITE)

    total = stats["total"]
    alert_count = stats["alert_count"]
    reg_count = stats["reg_count"]
    warning_count = stats["warning_count"]
    top_countries = stats["top_countries"]
    country_list = ", ".join(c for c, _ in top_countries[:5])

    # Big impact numbers
    boxes = [
        (f"{total:,}", "Suppliers Continuously\nMonitored"),
        (f"{reg_count}", "Regulations Tracked\nAcross Portfolio"),
        (f"{alert_count}", "Active Risk Alerts\nFlagged"),
        (f"{warning_count}", "Trade & Compliance\nWarnings Identified"),
    ]
    for i, (val, label) in enumerate(boxes):
        left = Inches(0.3 + i * 2.4)
        add_rect(slide, left, Inches(1.1), Inches(2.1), Inches(1.4),
                 FRDM_ACCENT if i == 0 else FRDM_BLUE)
        add_textbox(slide, val,
                    left + Inches(0.1), Inches(1.15), Inches(1.9), Inches(0.75),
                    font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, label,
                    left + Inches(0.05), Inches(1.85), Inches(2.0), Inches(0.6),
                    font_size=8, color=MID_GREY, align=PP_ALIGN.CENTER)

    # Narrative bullets
    bullets = [
        f"Monitoring {total:,} suppliers across {len(top_countries)} countries including {country_list}",
        f"Tracking {reg_count} unique regulatory requirements that affect {company_name}'s supply chain",
        f"Identified {warning_count} suppliers with active warnings including forced labor, trade restrictions, and tariff risks",
        f"Flagged {alert_count} suppliers with live risk alerts requiring review",
        "Continuously updating risk scores as geopolitical events, sanctions, and regulations change",
        "Providing transparency scoring to drive supplier engagement and disclosure improvements",
    ]

    add_textbox(slide, "FRDM's Ongoing Work for " + company_name,
                Inches(0.5), Inches(2.75), Inches(9), Inches(0.4),
                font_size=12, bold=True, color=RGBColor(0xCC, 0xD6, 0xFF))
    bullet_list(slide, bullets,
                Inches(0.5), Inches(3.2), Inches(9), Inches(3.2),
                font_size=10, color=RGBColor(0xDD, 0xE4, 0xFF))

    add_textbox(slide, "Prepared by FRDM  ·  Confidential  ·  " + date.today().strftime("%B %Y"),
                Inches(0.5), Inches(6.8), Inches(9), Inches(0.35),
                font_size=8, color=MID_GREY, align=PP_ALIGN.CENTER)


# ── Main build function ───────────────────────────────────────────────────────

def build_qbr(df: pd.DataFrame, company_name: str, quarter: str) -> bytes:
    # ── Compute stats ──────────────────────────────────────────────────────────
    df["_risk_score"]  = df["topline_risk_score"].apply(safe_float)
    df["_trans_score"] = df["transparency_score"].apply(safe_float)
    df["_risk_label"]  = df["_risk_score"].apply(risk_label)
    df["_trans_label"] = df["_trans_score"].apply(transparency_label)
    df["_alert_count"] = pd.to_numeric(df.get("all_alerts_count", 0), errors="coerce").fillna(0)

    total = len(df)
    avg_risk  = df["_risk_score"].dropna().mean()
    avg_trans = df["_trans_score"].dropna().mean()

    risk_dist  = df["_risk_label"].value_counts().to_dict()
    trans_dist = df["_trans_label"].value_counts().to_dict()

    alert_count   = int((df["_alert_count"] > 0).sum())
    warning_count = int(df["Warnings"].fillna("").astype(str).str.strip().ne("").sum())
    doc_count     = int(df["Documents"].fillna("").astype(str).str.strip().ne("").sum())

    all_warnings  = Counter()
    all_docs      = Counter()
    all_prod_regs = Counter()
    all_ind_regs  = Counter()
    all_countries = Counter()

    for _, row in df.iterrows():
        for w in parse_pipe(row.get("Warnings", "")):      all_warnings[w] += 1
        for d in parse_pipe(row.get("Documents", "")):     all_docs[d] += 1
        for p in parse_pipe(row.get("Product Regulations", "")): all_prod_regs[p] += 1
        for i in parse_pipe(row.get("Industry Regulations", "")): all_ind_regs[i] += 1
        for c in parse_pipe(row.get("Countries", "")):     all_countries[c] += 1

    reg_count = len(set(list(all_prod_regs.keys()) + list(all_ind_regs.keys())))

    stats = dict(
        total=total,
        avg_risk=avg_risk,
        avg_trans=avg_trans,
        risk_dist=risk_dist,
        trans_dist=trans_dist,
        alert_count=alert_count,
        warning_count=warning_count,
        doc_count=doc_count,
        reg_count=reg_count,
        top_warnings=all_warnings.most_common(10),
        top_docs=all_docs.most_common(10),
        top_prod_regs=all_prod_regs.most_common(10),
        top_ind_regs=all_ind_regs.most_common(10),
        top_countries=all_countries.most_common(10),
    )

    top_risky = df.nlargest(10, "_risk_score").to_dict("records")
    warning_suppliers = (
        df[df["Warnings"].fillna("").astype(str).str.strip() != ""]
        .nlargest(8, "_risk_score")
        .to_dict("records")
    )

    # ── Build PPTX ────────────────────────────────────────────────────────────
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_title(prs, company_name, quarter)
    slide_executive_summary(prs, stats)
    slide_risk_landscape(prs, stats, top_risky)
    slide_warnings(prs, stats, warning_suppliers)
    slide_regulations(prs, stats)
    slide_documents(prs, stats)
    slide_frdm_value(prs, stats, company_name)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


# ── Streamlit UI ──────────────────────────────────────────────────────────────

st.markdown(
    """
    <div style="padding:1.5rem 0 0.5rem 0;">
        <h1 style="font-size:2rem; font-weight:700; color:#1a1a2e; margin:0;">
            📋 QBR Generator
        </h1>
        <p style="color:#555; margin-top:0.25rem; font-size:1rem;">
            Upload a FRDM supplier risk CSV to generate a customer-ready PowerPoint QBR deck.
        </p>
    </div>
    <hr style="border:none; border-top:2px solid #e8eaf0; margin-bottom:1.5rem;">
    """,
    unsafe_allow_html=True,
)

col_left, col_right = st.columns([2, 1], gap="large")

with col_left:
    uploaded = st.file_uploader(
        "Upload supplier risk CSV",
        type=["csv"],
        help="Export from the FRDM platform — supplier-level report with risk scores, warnings, regulations, etc.",
    )

with col_right:
    company_name_input = st.text_input(
        "Customer name (optional)",
        placeholder="Leave blank to use Company ID from CSV",
        help="This will appear on the title slide. If left blank, the company ID from the CSV is used.",
    )
    quarter_input = st.text_input(
        "Quarter / period label",
        value=f"Q{((date.today().month - 1) // 3) + 1} {date.today().year}",
    )

if uploaded:
    try:
        df = pd.read_csv(uploaded)
        required = {"topline_risk_score", "transparency_score"}
        missing = required - set(df.columns)
        if missing:
            st.error(f"CSV is missing required columns: {missing}")
        else:
            company_id = str(df["id"].iloc[0]) if "id" in df.columns else "Unknown"
            display_name = company_name_input.strip() if company_name_input.strip() else f"Company {company_id}"

            st.success(
                f"Loaded **{len(df):,}** suppliers · Company ID: **{company_id}** · "
                f"Will use name: **{display_name}**"
            )

            # Preview stats
            with st.expander("Preview data summary", expanded=False):
                c1, c2, c3, c4 = st.columns(4)
                c1.metric("Total Suppliers", f"{len(df):,}")
                avg_r = pd.to_numeric(df["topline_risk_score"], errors="coerce").mean()
                c2.metric("Avg Risk Score", f"{avg_r:.3f}" if not pd.isna(avg_r) else "N/A")
                avg_t = pd.to_numeric(df["transparency_score"], errors="coerce").mean()
                c3.metric("Avg Transparency", f"{avg_t:.3f}" if not pd.isna(avg_t) else "N/A")
                alert_n = pd.to_numeric(df.get("all_alerts_count", pd.Series()), errors="coerce").gt(0).sum()
                c4.metric("Suppliers w/ Alerts", int(alert_n))

            if st.button("Generate QBR Deck", type="primary", use_container_width=True):
                with st.spinner("Building PowerPoint..."):
                    try:
                        pptx_bytes = build_qbr(df, display_name, quarter_input)
                        safe_name = display_name.replace(" ", "_")
                        fname = f"FRDM_QBR_{safe_name}_{quarter_input.replace(' ', '_')}.pptx"
                        st.download_button(
                            label="Download QBR (.pptx)",
                            data=pptx_bytes,
                            file_name=fname,
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            use_container_width=True,
                        )
                        st.success("QBR deck ready! Click above to download.")
                    except Exception as e:
                        st.error(f"Error building deck: {e}")
                        raise

    except Exception as e:
        st.error(f"Could not read CSV: {e}")
else:
    st.info("Upload a CSV to get started.")
